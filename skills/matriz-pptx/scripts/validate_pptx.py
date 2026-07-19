#!/usr/bin/env python3
"""Detect obvious off-canvas shapes and missing visible content in a PPTX."""
import argparse, json
from pathlib import Path
from pptx import Presentation

def main():
    parser = argparse.ArgumentParser(); parser.add_argument("input", type=Path); parser.add_argument("--out", type=Path)
    args = parser.parse_args(); deck = Presentation(args.input); issues = []
    for number, slide in enumerate(deck.slides, 1):
        if not slide.shapes: issues.append({"slide": number, "issue": "slide has no shapes"})
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > deck.slide_width or shape.top + shape.height > deck.slide_height:
                issues.append({"slide": number, "shape": shape.name, "issue": "shape extends beyond slide canvas"})
    report = json.dumps({"valid": not issues, "issues": issues}, indent=2) + "\n"
    if args.out: args.out.write_text(report, encoding="utf-8")
    else: print(report, end="")
    if issues: raise SystemExit(1)

if __name__ == "__main__": main()

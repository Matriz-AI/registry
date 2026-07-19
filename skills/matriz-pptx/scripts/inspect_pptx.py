#!/usr/bin/env python3
"""Report slide and asset structure without modifying a PPTX."""
import argparse, json
from pathlib import Path
from pptx import Presentation

def main():
    parser = argparse.ArgumentParser(); parser.add_argument("input", type=Path); parser.add_argument("--out", type=Path)
    args = parser.parse_args(); deck = Presentation(args.input); slides = []
    for number, slide in enumerate(deck.slides, 1):
        titles = [shape.text.strip() for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        slides.append({"number": number, "shapes": len(slide.shapes), "title": titles[0] if titles else None, "notes": bool(slide.notes_slide.notes_text_frame.text.strip())})
    text = json.dumps({"file": args.input.name, "size": [deck.slide_width, deck.slide_height], "slides": slides}, indent=2) + "\n"
    if args.out: args.out.write_text(text, encoding="utf-8")
    else: print(text, end="")

if __name__ == "__main__": main()

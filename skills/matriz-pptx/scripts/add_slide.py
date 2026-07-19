#!/usr/bin/env python3
"""Append an editable title-and-body slide while preserving the source deck."""
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

def main():
    parser = argparse.ArgumentParser(); parser.add_argument("input", type=Path); parser.add_argument("output", type=Path)
    parser.add_argument("--title", required=True); parser.add_argument("--body", default="")
    args = parser.parse_args(); deck = Presentation(args.input); slide = deck.slides.add_slide(deck.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(.7), Inches(.45), deck.slide_width - Inches(1.4), Inches(.8)).text_frame
    title.text = args.title; title.paragraphs[0].font.size = Pt(30); title.paragraphs[0].font.bold = True
    body = slide.shapes.add_textbox(Inches(.8), Inches(1.55), deck.slide_width - Inches(1.6), deck.slide_height - Inches(2.1)).text_frame
    body.text = args.body
    for paragraph in body.paragraphs: paragraph.font.size = Pt(18)
    deck.save(args.output)

if __name__ == "__main__": main()
